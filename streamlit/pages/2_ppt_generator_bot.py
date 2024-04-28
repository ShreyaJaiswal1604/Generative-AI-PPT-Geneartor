import streamlit as st
from openai import Client
import time
from PIL import Image
from openai import OpenAI
import io
import os
import base64

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

openai_api_key=os.getenv('OPENAI_API_KEY')
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))


def sidebar_bg(side_bg):

   side_bg_ext = 'png'

   st.markdown(
      f"""
      <style>
      [data-testid="stSidebar"] > div:first-child {{
          background: url(data:image/{side_bg_ext};base64,{base64.b64encode(open(side_bg, "rb").read()).decode()});
      }}
      </style>
      """,
      unsafe_allow_html=True,
      )
   
side_bg_ext = '../media/images/img-07.jpeg'

sidebar_bg(side_bg_ext)

file_list = os.listdir("./tmp")
filtered_list = [filename for filename in file_list if filename.endswith('.json')]

file_selected = st.selectbox(
    'Please select a file from the list. If you cannot find the file head over to page 1 to upload it', 
    filtered_list
    )


global image_path

def get_response(thread):
    return client.beta.threads.messages.list(thread_id=thread.id)

def convert_file_to_png(file_id, write_path):
    data = client.files.content(file_id)
    data_bytes = data.read()
    with open(write_path, "wb") as file:
        file.write(data_bytes)

st.title("PPT Generator Bot")
insight_query_one = st.text_area(label = "Enter Business Query", 
                                 value= "Calculate profit (revenue minus cost) by quarter and year, and visualize as a line plot across the distribution channels, where the colors of the lines are green, light red, and light blue")

file = client.files.create(
    file=open(f"./tmp/{file_selected}", "rb"),
    purpose='assistants'
)

# Create assistant with the file
assistant = client.beta.assistants.create(
    name="Data visualizer",
    description="Your assistant description here...",
    model="gpt-4-1106-preview",
    tools=[{"type": "code_interpreter"}],
    tool_resources={
        "code_interpreter": {
            "file_ids": [file.id]
        }
    }
)


if st.button("Generate Query Result"):
    # Create file in OpenAI


    # Create a thread to communicate with the assistant
    thread = client.beta.threads.create(
        messages=[
            {
                "role": "user",
                "content": f"{insight_query_one}",
                "attachments": [
                    {
                        "file_id": file.id,
                        "tools": [{"type": "code_interpreter"}]
                    }
                ]
            }
        ]
    )

    # Run the assistant
    run = client.beta.threads.runs.create(
        thread_id=thread.id,
        assistant_id=assistant.id
    )

    # Poll assistant to check if process is completed
    with st.spinner(text='Assistant still working...'):
        while True:
            messages = client.beta.threads.messages.list(thread_id=thread.id)
            try:
                # Check if image has been created
                messages.data[0].content[0].image_file
                # Print message once plot is created
                st.success('Plot created!')
                break
            except:
                # Print message while assistant is still working
                #st.warning('Assistant still working...')
                time.sleep(10)
                continue


    # Display the output image to the user
    for message in messages.data:
        if message.content[0].type == "image_file":
            file_id = message.content[0].image_file.file_id
            plot_summary = message.content[1].text.value
            image_bytes = client.files.content(file_id).read()
            image = Image.open(io.BytesIO(image_bytes))
            st.image(image, caption=plot_summary, use_column_width=True)

            # Save the image
            image_path = f"./tmp/{file_selected}_plot.png"
            # image.save(image_path)
            with open(image_path, 'wb') as f:
                image.save(f)
            st.write(f"Image saved as {image_path}") # TODO: Piyush fix this location
    # Add button to generate final PPT


st.title("Company Summary Input")
# Get company summary input from user
company_summary = st.text_input(label = "Enter company summary:", value = "Elevance Health - Medical Insurance Company")
# Get inputs from user
title_text = st.text_input(label = "Enter title text:", value = "Quarterly Report 2024")
subtitle_text = st.text_input(label = "Enter subtitle text:", value = "Quarterly Report 2024")
output_file_name = st.text_input(label = "Enter output filename:", value = "financial-report-quarterly-2024")
business_summary = st.text_input(label = "Enter business summary:", value = "Medical Insurance Company insurance company that manufactures and sells health policies, graphics cards and other essential medical healthcare")


if st.button(label = "Generate Final PPT"
            #  disabled = True
             ):
    #Upload
    with st.spinner(text='Assistant working on PPTX...'):
        plot_file = client.files.create(
            file=open(f"./tmp/{file_selected}_plot.png", "rb"),
            purpose='assistants'
        )

        response = client.images.generate(
        model='dall-e-3',
        prompt=f"given this company summary {company_summary}, create an inspirational \
            photo showing the growth and path forward. This will be used at a quarterly\
            financial planning meeting",
            size="1024x1024",
            quality="standard",
            n=1
        )
        image_url = response.data[0].url

        dalle_img_path = f"./tmp/{file_selected}_cover.png"
        img = requests.get(image_url)

        #Save locally
        with open(dalle_img_path,'wb') as file:
            file.write(img.content)

        #Upload
        dalle_file_img = client.files.create(
        file=open(dalle_img_path, "rb"),
        purpose='assistants'
        )
        
        # Code templates for creating PPT slides
        title_template = f"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor

        # Create a new presentation object
        prs = Presentation()

        # Add a blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set the background color of the slide to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Add image to the left side of the slide with a margin at the top and bottom
        left = Inches(0)
        top = Inches(0)
        height = prs.slide_height
        width = prs.slide_width * 3/5
        pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

        # Add title text box positioned higher
        left = prs.slide_width * 3/5
        top = Inches(2)
        width = prs.slide_width * 2/5
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_p = title_frame.add_paragraph()
        title_p.text = "{title_text}"
        title_p.font.bold = True
        title_p.font.size = Pt(38)
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Add subtitle text box
        left = prs.slide_width * 3/5
        top = Inches(3)
        width = prs.slide_width * 2/5
        height = Inches(1)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.add_paragraph()
        subtitle_p.text = "{subtitle_text}"
        subtitle_p.font.size = Pt(22)
        subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Wrap text in title and subtitle text boxes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
        """

        data_vis_template = f"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor

        # Create a new presentation object
        prs = Presentation()

        # Add a blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set the background color of the slide to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Define placeholders
        image_path = data_vis_img
        title_text = "{title_text}"
        bullet_points = "{business_summary}"
        # Add image placeholder on the left side of the slide
        left = Inches(0.2)
        top = Inches(1.8)
        height = prs.slide_height - Inches(3)
        width = prs.slide_width * 3/5
        pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)

        # Add title text spanning the whole width
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.margin_top = Inches(0.1)
        title_p = title_frame.add_paragraph()
        title_p.text = title_text
        title_p.font.bold = True
        title_p.font.size = Pt(28)
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Add hardcoded "Key Insights" text and bullet points
        left = prs.slide_width * 2/3
        top = Inches(1.5)
        width = prs.slide_width * 1/3
        height = Inches(4.5)
        insights_box = slide.shapes.add_textbox(left, top, width, height)
        insights_frame = insights_box.text_frame
        insights_p = insights_frame.add_paragraph()
        insights_p.text = "Key Insights:"
        insights_p.font.bold = True
        insights_p.font.size = Pt(24)
        insights_p.font.color.rgb = RGBColor(0, 128, 100)
        insights_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        insights_frame.add_paragraph()

        bullet_p = insights_frame.add_paragraph()
        bullet_p.text = bullet_points
        bullet_p.font.size = Pt(12)
        bullet_p.font.color.rgb = RGBColor(255, 255, 255)
        bullet_p.line_spacing = 1.5

        # Wrap text in title and bullet points text boxes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
        """

        thank_you_template = f"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor

        # Create a new presentation object
        prs = Presentation()

        # Add a blank slide layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set the background color of the slide to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Add title text box positioned in the middle of the slide
        left = Inches(1)
        top = Inches(2)
        width = prs.slide_width - Inches(2)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_p = title_frame.add_paragraph()
        title_p.text = "Thank You!"
        title_p.font.bold = True
        title_p.font.size = Pt(40)
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        # Wrap text in title text box
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
        """

        promptone = f"Use the included code template to create a PPTX slide that follows the template format, but uses the image, company name/title, and document name/subtitle included:\
            {title_template}. IMPORTANT: Use the first image file included in this message attachment as the image_path in this first slide, and use the Company Name {title_text} as the title_text variable, and \
            use the subtitle_text {subtitle_text} a the subtitle_text variable. \
                NEXT, create a SECOND slide using the following code template: {data_vis_template} to create a PPTX slide that follows the template format, but uses the company name/title, and document name/subtitle included:\
            {data_vis_template}. IMPORTANT: Use the line plot image, that is the second attached image in this message, that you created earlier in the thread as the data_vis_img image, and use the data visualization title that you created earlier for the variable title_text, and\
            the bullet points of insights you created earlier for the bullet_points variable.\
            NEXT, create a THIRD slide using the following code template: {thank_you_template} to create a PPTX slide that follows the template format.\
            Output these THREE SLIDES as a .pptx file. Make sure the output is three slides, with each slide matching the respective template given in this message. IMPORTANT: if the image files are not in the proper format, convert them  into proper format and then use."
        
        # Create a thread and attach the files to the message
        thread = client.beta.threads.create(
            messages=[
                {
                    "role": "user",
                    "content": f"{promptone}",
                    "attachments": [
                        {"file_id": dalle_file_img.id, "tools": [{"type": "code_interpreter"}]},
                        {"file_id": plot_file.id, "tools": [{"type": "code_interpreter"}]}
                        
                    ]
                }
            ]
        )

        # The thread now has a vector store with that file in its tool resources.
        print(thread.tool_resources.code_interpreter)
        # Use the create and poll SDK helper to create a run and poll the status of
        # the run until it's in a terminal state.

        print("running client.....")
        run = client.beta.threads.runs.create_and_poll(
            thread_id=thread.id, assistant_id=assistant.id
        )

        messages = list(client.beta.threads.messages.list(thread_id=thread.id, run_id=run.id))

        for message in messages:
            print(message.content[0].text.value)

        #May take 1-3 mins
        while True:
            try:
                response = get_response(thread)
                pptx_id = response.data[0].content[0].text.annotations[0].file_path.file_id
                print("Successfully retrieved pptx_id:", pptx_id)
                break
            except Exception as e:
                print("Assistant still working on PPTX...")
                time.sleep(10)

        pptx_id = response.data[0].content[0].text.annotations[0].file_path.file_id
        ppt_file= client.files.content(pptx_id)
        file_obj = io.BytesIO(ppt_file.read())

        with open(f"./tmp/{title_text}.pptx", "wb") as f:
            f.write(file_obj.getbuffer())

        btn = st.download_button(
            label="Download pptx",
            data=file_obj,
            file_name=f"{title_text}.pptx",
            mime="application/octet-stream"
        )

        client.files.delete(pptx_id)
        client.beta.assistants.delete(assistant.id)
        print("successfully deleted assistant")

    #generate_ppt(title_text,subtitle_text,output_file_name,business_summary)

    # st.write("Generating final PPT...")
    # # Redirect to the final PPT generation page
    # st.experimental_reroute("/generator-final-ppt")